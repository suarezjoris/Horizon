import sys
import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

def get_branding(vault_path):
    branding_path = os.path.join(vault_path, "assets", "branding.json")
    if os.path.exists(branding_path):
        try:
            with open(branding_path, 'r') as f:
                return json.load(f)
        except:
            pass
    return None

def apply_clay_style(slide, x, y, w, h, palette):
    """Applies Claymorphism style to a shape with a subtle shadow."""
    # 1. Shadow/Glow (Bottom layer)
    shadow = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.05), y + Inches(0.05), w, h
    )
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = RGBColor(0, 0, 0)
    shadow.fill.transparency = 0.6
    shadow.line.fill.background() # No outline

    # 2. The main 'Clay' body (Top layer)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h
    )
    fill = shape.fill
    fill.solid()
    r, g, b = palette.get("surface", [40, 40, 40])
    fill.fore_color.rgb = RGBColor(r, g, b)
    
    # Inner highlights simulation via line
    line = shape.line
    line.color.rgb = RGBColor(r+30 if r<225 else 255, g+30 if g<225 else 255, b+30 if b<225 else 255)
    line.width = Pt(2)
    return shape

def generate_pptx(data):
    vault_path = os.path.dirname(os.path.dirname(data.get("output_path")))
    branding = get_branding(vault_path)
    palette = branding.get("palette", {}) if branding else {}
    
    template_path = data.get("template_path")
    prs = Presentation(template_path) if template_path and os.path.exists(template_path) else Presentation()

    # Colors
    bg_color = RGBColor(*palette.get("background", [26, 26, 26]))
    accent_color = RGBColor(*palette.get("accent", [212, 175, 55]))
    text_color = RGBColor(*palette.get("text_main", [255, 255, 255]))

    # 1. Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg_color

    # Center box
    title_box = apply_clay_style(slide, Inches(1), Inches(2), Inches(8), Inches(3.5), palette)
    
    tf = title_box.text_frame
    tf.text = data.get("title", "Presentation")
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = accent_color
    p.alignment = PP_ALIGN.CENTER

    # 2. Content Slides
    for slide_data in data.get("slides", []):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg_color

        # Decorative side accent (Gold bar)
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.1), Inches(7.5)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = accent_color
        accent_bar.line.fill.background()

        # Header bar
        header = apply_clay_style(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.8), palette)
        header.text_frame.text = slide_data.get("title", "")
        p = header.text_frame.paragraphs[0]
        p.font.color.rgb = accent_color
        p.font.bold = True
        p.font.size = Pt(28)
        p.alignment = PP_ALIGN.LEFT

        # Content Box
        body = apply_clay_style(slide, Inches(0.5), Inches(1.4), Inches(9), Inches(5.8), palette)
        tf = body.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = slide_data.get("intro", "")
        p.font.color.rgb = text_color
        p.font.size = Pt(20)
        p.font.italic = True
        
        for bullet in slide_data.get("bullets", []):
            bp = tf.add_paragraph()
            bp.text = f"  • {bullet}"
            bp.level = 0
            bp.font.color.rgb = RGBColor(*palette.get("text_dim", [180, 180, 180]))
            bp.font.size = Pt(18)
            bp.font.bold = False

    filename = data.get("filename", "presentation.pptx")
    output_path = data.get("output_path", filename)
    prs.save(output_path)
    print(output_path)

if __name__ == "__main__":
    if len(sys.argv) < 2:
        sys.exit(1)
    
    try:
        raw_data = sys.argv[1]
        data = json.loads(raw_data)
        generate_pptx(data)
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(1)
