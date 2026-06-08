import sys
import os
import docx
from pptx import Presentation

def extract_text(path):
    ext = os.path.splitext(path)[1].lower()
    
    if ext == ".docx":
        doc = docx.Document(path)
        return "\n".join([para.text for para in doc.paragraphs])
    
    elif ext == ".pptx":
        prs = Presentation(path)
        text_runs = [
            shape.text
            for slide in prs.slides
            for shape in slide.shapes
            if hasattr(shape, "text")
        ]
        return "\n".join(text_runs)
    
    elif ext in [".xlsx", ".xls"]:
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        text_runs = []
        for sheet in wb.worksheets:
            text_runs.append(f"--- Sheet: {sheet.title} ---")
            for row in sheet.iter_rows(values_only=True):
                # Filter out None values and join row data
                row_data = [str(cell) for cell in row if cell is not None]
                if row_data:
                    text_runs.append(" | ".join(row_data))
        wb.close()
        return "\n".join(text_runs)
    
    return ""

if __name__ == "__main__":
    if len(sys.argv) < 2:
        sys.exit(1)
    
    file_path = sys.argv[1]
    if not os.path.exists(file_path):
        print(f"File not found: {file_path}", file=sys.stderr)
        sys.exit(1)
        
    try:
        print(extract_text(file_path))
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(1)
